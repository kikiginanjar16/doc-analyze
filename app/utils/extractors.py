from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Optional, List, Tuple

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

from .config import SCAN_TEXT_THRESHOLD
from .openai_client import vision_page_to_markdown

@dataclass
class ExtractResult:
    combined_text: str
    meta: dict
    used_vision: bool = False

def _clean_text(s: str) -> str:
    s = s.replace("\x00", "")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def extract_txt_md(data: bytes, filename: str) -> ExtractResult:
    # decode with fallback
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1", errors="replace")
    text = _clean_text(text)
    meta = {"detected_type": "md" if filename.lower().endswith(".md") else "txt"}
    return ExtractResult(combined_text=text, meta=meta, used_vision=False)

def extract_docx(data: bytes, max_sections: int = 500) -> ExtractResult:
    doc = DocxDocument(io.BytesIO(data))
    parts: List[str] = []
    # paragraphs
    for p in doc.paragraphs[:max_sections]:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    # tables (flatten)
    for ti, table in enumerate(doc.tables[:50]):
        parts.append(f"[Table {ti+1}]")
        for row in table.rows[:200]:
            cells = [ (c.text or "").strip() for c in row.cells ]
            cells = [c for c in cells if c]
            if cells:
                parts.append(" | ".join(cells))
    text = _clean_text("\n".join(parts))
    meta = {"detected_type": "docx"}
    return ExtractResult(combined_text=text, meta=meta, used_vision=False)

def extract_pptx(data: bytes, max_slides: int = 200) -> ExtractResult:
    prs = Presentation(io.BytesIO(data))
    parts: List[str] = []
    slide_count = len(prs.slides)
    limit = min(slide_count, max_slides)
    for i in range(limit):
        slide = prs.slides[i]
        slide_parts: List[str] = []
        title = ""
        # gather shape texts
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if not t:
                continue
            if not title and shape.has_text_frame and shape.text_frame and shape.text_frame.text:
                # heuristic: first text could be title
                title = t.split("\n", 1)[0].strip()
            slide_parts.append(t)
        # speaker notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        parts.append(f"[Slide {i+1}] {title}".strip())
        if slide_parts:
            parts.extend(slide_parts)
        if notes:
            parts.append(f"[Slide {i+1} Notes] {notes}")
        parts.append("")  # spacing
    text = _clean_text("\n".join(parts))
    meta = {"detected_type": "pptx", "slides": slide_count}
    return ExtractResult(combined_text=text, meta=meta, used_vision=False)

def extract_pdf_text_or_vision(
    data: bytes,
    max_pages: int,
    vision_system_prompt: str,
    mime_for_images: str = "image/png",
) -> ExtractResult:
    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as e:
        raise ValueError("Invalid or unreadable PDF file.") from e
    page_count = doc.page_count
    limit = min(page_count, max_pages)
    per_page_text: List[str] = []
    total_chars = 0

    # text extraction pass (limited)
    for i in range(limit):
        page = doc.load_page(i)
        t = page.get_text("text") or ""
        t = _clean_text(t)
        if t:
            per_page_text.append(f"[Page {i+1}]\n{t}\n")
            total_chars += len(t)

    # decide scanned
    scanned = total_chars < SCAN_TEXT_THRESHOLD
    used_vision = False
    if scanned:
        used_vision = True
        per_page_text = []
        for i in range(limit):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=200, alpha=False)
            img_bytes = pix.tobytes("png")
            hint = f"This is page {i+1} of a PDF document. Extract meaningful content and summarize."
            md = vision_page_to_markdown(img_bytes, mime_for_images, vision_system_prompt, hint=hint)
            per_page_text.append(f"## Page {i+1}\n{md}\n")
        combined = _clean_text("\n".join(per_page_text))
        meta = {"detected_type": "pdf", "pages": page_count, "used_vision_pages": limit, "scanned_detected": True}
        return ExtractResult(combined_text=combined, meta=meta, used_vision=True)

    combined = _clean_text("\n".join(per_page_text))
    meta = {"detected_type": "pdf", "pages": page_count, "scanned_detected": False}
    return ExtractResult(combined_text=combined, meta=meta, used_vision=False)
